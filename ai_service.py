"""
TeachFlowUZ — AI Service
Gemma-4 orqali rasm tahlili, DeepSeek-R1 orqali baho va feedback
"""

import os
import json
import base64
import traceback
from io import BytesIO

from huggingface_hub import InferenceClient
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image

# ── API Configuration ─────────────────────────────────────────────
HF_API_KEY = os.getenv("HUGGINGFACE_API_KEY", "your_token_here_if_testing_locally")
IMAGE_MODEL = "google/gemma-4-31B-it:novita"
TEXT_MODEL = "deepseek-ai/DeepSeek-R1:novita"

client = InferenceClient(api_key=HF_API_KEY)


# ── File Content Extraction ──────────────────────────────────────

def extract_pdf_content(filepath):
    """PDF fayldan matn va rasmlarni ajratish"""
    result = {"text": "", "images": []}
    try:
        reader = PdfReader(filepath)
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                result["text"] += f"\n--- Sahifa {page_num + 1} ---\n{text}"

            # PDF dan rasmlarni ajratish
            if hasattr(page, 'images'):
                for img_idx, image in enumerate(page.images):
                    try:
                        img_data = image.data
                        img_base64 = base64.b64encode(img_data).decode('utf-8')
                        result["images"].append({
                            "page": page_num + 1,
                            "index": img_idx,
                            "data": img_base64,
                            "name": image.name
                        })
                    except Exception:
                        pass
    except Exception as e:
        result["error"] = str(e)
    return result


def extract_docx_content(filepath):
    """DOCX fayldan matn va rasmlarni ajratish"""
    result = {"text": "", "images": []}
    try:
        doc = DocxDocument(filepath)

        # Matnni ajratish
        for para_idx, para in enumerate(doc.paragraphs):
            if para.text.strip():
                style_name = para.style.name if para.style else "Normal"
                if "Heading" in style_name:
                    result["text"] += f"\n\n## {para.text}\n"
                else:
                    result["text"] += f"{para.text}\n"

        # Rasmlarni ajratish
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_data = rel.target_part.blob
                    img_base64 = base64.b64encode(img_data).decode('utf-8')
                    result["images"].append({
                        "data": img_base64,
                        "name": rel.target_ref
                    })
                except Exception:
                    pass
    except Exception as e:
        result["error"] = str(e)
    return result


def extract_pptx_content(filepath):
    """PPTX fayldan matn va rasmlarni ajratish"""
    result = {"text": "", "images": [], "slide_count": 0}
    try:
        prs = Presentation(filepath)
        result["slide_count"] = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides):
            result["text"] += f"\n\n=== Slayd {slide_idx + 1} ===\n"

            for shape in slide.shapes:
                # Matn ajratish
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            result["text"] += f"{paragraph.text}\n"

                # Rasm ajratish
                if shape.shape_type == 13:  # Picture type
                    try:
                        img_data = shape.image.blob
                        img_base64 = base64.b64encode(img_data).decode('utf-8')
                        result["images"].append({
                            "slide": slide_idx + 1,
                            "data": img_base64,
                            "name": f"slide_{slide_idx + 1}_image"
                        })
                    except Exception:
                        pass
    except Exception as e:
        result["error"] = str(e)
    return result


def extract_video_frames(filepath, max_frames=5):
    """Videodan kadrlarni ajratish"""
    result = {"frames": [], "duration": 0}
    try:
        import cv2
        cap = cv2.VideoCapture(filepath)
        if not cap.isOpened():
            result["error"] = "Video faylni ochib bo'lmadi"
            return result

        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = cap.get(cv2.CAP_PROP_FPS)
        if fps > 0:
            result["duration"] = total_frames / fps

        # Teng intervallarda kadrlarni ajratish
        frame_indices = [int(i * total_frames / (max_frames + 1)) for i in range(1, max_frames + 1)]

        for idx in frame_indices:
            cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
            ret, frame = cap.read()
            if ret:
                # Kadrni Base64 ga o'girish
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                img = Image.fromarray(frame_rgb)
                img.thumbnail((800, 600))
                buffer = BytesIO()
                img.save(buffer, format='JPEG', quality=75)
                img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
                result["frames"].append({
                    "timestamp": round(idx / fps, 1) if fps > 0 else 0,
                    "data": img_base64
                })

        cap.release()
    except ImportError:
        result["error"] = "opencv-python o'rnatilmagan"
    except Exception as e:
        result["error"] = str(e)
    return result


# ── AI Analysis Functions ────────────────────────────────────────

def analyze_image_with_ai(image_base64, context=""):
    """Gemma-4 orqali rasmni tahlil qilish"""
    try:
        prompt = "Bu rasmni batafsil tavsiflang. Rasmda nima tasvirlangan, qanday matn bor, qanday diagramma yoki grafik bor — hammasini batafsil yozing."
        if context:
            prompt += f" Kontekst: {context}"

        data_url = f"data:image/jpeg;base64,{image_base64}"

        completion = client.chat.completions.create(
            model=IMAGE_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": data_url}}
                    ]
                }
            ],
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"[Rasmni tahlil qilib bo'lmadi: {str(e)}]"


def grade_with_ai(full_prompt):
    """DeepSeek-R1 orqali baho va feedback olish"""
    try:
        completion = client.chat.completions.create(
            model=TEXT_MODEL,
            messages=[
                {
                    "role": "user",
                    "content": full_prompt
                }
            ],
        )
        return completion.choices[0].message.content
    except Exception as e:
        return json.dumps({
            "error": str(e),
            "score": 0,
            "feedback": "AI xizmatiga ulanib bo'lmadi. Iltimos, keyinroq urinib ko'ring."
        })


# ── Main Evaluation Pipeline ────────────────────────────────────

def evaluate_submission(filepath, filename, file_type, assignment_title, assignment_description):
    """
    To'liq tekshirish jarayoni:
    1. Fayldan matn va rasmlarni ajratish
    2. Rasmlarni Gemma-4 orqali tahlil qilish
    3. Hammasini birlashtirib DeepSeek-R1 ga yuborish
    4. Baho va feedback olish
    """
    result = {
        "score": 0,
        "feedback": "",
        "strengths": "",
        "weaknesses": "",
        "recommendations": "",
        "detailed_analysis": "",
        "status": "processing"
    }

    try:
        # 1-QADAM: Fayldan kontent ajratish
        content_data = {"text": "", "images": []}

        if file_type == 'pdf':
            content_data = extract_pdf_content(filepath)
        elif file_type == 'docx':
            content_data = extract_docx_content(filepath)
        elif file_type == 'pptx':
            content_data = extract_pptx_content(filepath)
        elif file_type in ('mp4', 'avi', 'mov', 'mkv'):
            video_data = extract_video_frames(filepath)
            content_data["text"] = f"[Video darslik: {filename}, Davomiyligi: {video_data.get('duration', 0):.1f} soniya]"
            for frame in video_data.get("frames", []):
                content_data["images"].append(frame)

        # 2-QADAM: Rasmlarni AI orqali tahlil qilish
        image_descriptions = []
        images_to_analyze = content_data.get("images", [])

        for i, img_info in enumerate(images_to_analyze[:10]):  # Max 10 ta rasm
            img_data = img_info.get("data", "")
            if img_data:
                location = ""
                if "page" in img_info:
                    location = f"Sahifa {img_info['page']}"
                elif "slide" in img_info:
                    location = f"Slayd {img_info['slide']}"
                elif "timestamp" in img_info:
                    location = f"Vaqt: {img_info['timestamp']}s"

                description = analyze_image_with_ai(
                    img_data,
                    context=f"Bu pedagogik material: '{assignment_title}'. {location}"
                )
                image_descriptions.append(f"[Rasm {i+1} ({location})]:\n{description}")

        # 3-QADAM: To'liq prompt tuzish
        full_prompt = f"""Sen pedagogika bo'yicha tajribali ekspert-o'qituvchisan. Talabaning topshirgan ishini batafsil tekshir va baho ber.

=== TOPSHIRIQ MA'LUMOTLARI ===
Topshiriq nomi: {assignment_title}
Topshiriq tavsifi: {assignment_description}
Yuborilgan fayl: {filename} ({file_type.upper()} format)

=== FAYLNING MATN KONTENTI ===
{content_data.get('text', '[Matn topilmadi]')}

"""
        if image_descriptions:
            full_prompt += "=== FAYLDAGI RASMLAR TAHLILI ===\n"
            full_prompt += "\n\n".join(image_descriptions)
            full_prompt += "\n\n"

        if file_type == 'pptx':
            full_prompt += f"Slaydlar soni: {content_data.get('slide_count', 'Noma\'lum')}\n\n"

        full_prompt += """=== BAHOLASH MEZONLARI ===
Quyidagi mezonlar bo'yicha baholab, har biriga 0-20 ball ber:
1. **Mavzuning to'liq yoritilganligi** (0-20): Mavzu qanchalik chuqur va to'liq yoritilgan?
2. **Ilmiy asoslanganlik** (0-20): Pedagogik nazariyalar, manbalar, ilmiy maqolalarga asoslangan-mi?
3. **Tuzilma va mantiqiylik** (0-20): Material qanchalik tartibli, mantiqiy ketma-ketlikda tuzilgan?
4. **Vizual dizayn va rasmiylashuv** (0-20): Slaydlar/hujjat dizayni, rasmlar, jadvallar sifati qanday?
5. **Amaliy qo'llanilishi** (0-20): Material amaliyotda qo'llasa bo'ladigan darajada tayyorlanganmi?

=== JAVOB FORMATI ===
Javobingni quyidagi aniq formatda ber (boshqa hech narsa qo'shma):

BAHO: [umumiy ball 0-100]

UMUMIY BAHOLASH:
[2-3 jumlada umumiy baho]

KUCHLI TOMONLAR:
• [1-chi kuchli tomon]
• [2-chi kuchli tomon]
• [3-chi kuchli tomon]

ZAIF TOMONLAR:
• [1-chi zaif tomon]
• [2-chi zaif tomon]
• [3-chi zaif tomon]

MEZONLAR BO'YICHA:
1. Mavzuning to'liq yoritilganligi: [ball]/20 — [izoh]
2. Ilmiy asoslanganlik: [ball]/20 — [izoh]
3. Tuzilma va mantiqiylik: [ball]/20 — [izoh]
4. Vizual dizayn va rasmiylashuv: [ball]/20 — [izoh]
5. Amaliy qo'llanilishi: [ball]/20 — [izoh]

TAVSIYALAR:
• [1-chi tavsiya]
• [2-chi tavsiya]
• [3-chi tavsiya]
• [4-chi tavsiya]
"""

        # 4-QADAM: AI dan baho olish
        ai_response = grade_with_ai(full_prompt)

        # 5-QADAM: Javobni parse qilish
        parsed = parse_ai_response(ai_response)
        result.update(parsed)
        result["status"] = "completed"

    except Exception as e:
        result["status"] = "error"
        result["feedback"] = f"Tekshirishda xatolik yuz berdi: {str(e)}\n{traceback.format_exc()}"

    return result


def parse_ai_response(response):
    """AI javobini parse qilish"""
    result = {
        "score": 0,
        "feedback": "",
        "strengths": "",
        "weaknesses": "",
        "recommendations": "",
        "detailed_analysis": response
    }

    try:
        lines = response.strip().split('\n')

        current_section = None
        sections = {
            "feedback": [],
            "strengths": [],
            "weaknesses": [],
            "recommendations": [],
        }

        for line in lines:
            line_stripped = line.strip()

            # Bahoni ajratish
            if line_stripped.startswith("BAHO:"):
                try:
                    score_text = line_stripped.replace("BAHO:", "").strip()
                    # Faqat raqamlarni olish
                    score_num = ''.join(c for c in score_text.split('/')[0] if c.isdigit() or c == '.')
                    if score_num:
                        result["score"] = min(100, max(0, float(score_num)))
                except (ValueError, IndexError):
                    pass

            elif "UMUMIY BAHOLASH:" in line_stripped:
                current_section = "feedback"
            elif "KUCHLI TOMONLAR:" in line_stripped:
                current_section = "strengths"
            elif "ZAIF TOMONLAR:" in line_stripped:
                current_section = "weaknesses"
            elif "MEZONLAR BO'YICHA:" in line_stripped:
                current_section = "feedback"
            elif "TAVSIYALAR:" in line_stripped:
                current_section = "recommendations"
            elif current_section and line_stripped:
                sections[current_section].append(line_stripped)

        result["feedback"] = '\n'.join(sections["feedback"])
        result["strengths"] = '\n'.join(sections["strengths"])
        result["weaknesses"] = '\n'.join(sections["weaknesses"])
        result["recommendations"] = '\n'.join(sections["recommendations"])

    except Exception:
        result["feedback"] = response

    return result
